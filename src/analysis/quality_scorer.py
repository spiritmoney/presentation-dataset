"""Analyze presentation files for graphics density and text heaviness."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.analysis.visual_clarity import analyze_clarity

VISUAL_SHAPE_TYPES = {
    MSO_SHAPE_TYPE.PICTURE,
    MSO_SHAPE_TYPE.CHART,
    MSO_SHAPE_TYPE.TABLE,
    MSO_SHAPE_TYPE.GROUP,
    MSO_SHAPE_TYPE.DIAGRAM,
    MSO_SHAPE_TYPE.CANVAS,
}


@dataclass
class QualityAnalysis:
    quality: float
    graphics_density: float
    text_density: float
    clarity: float
    modernity: float
    slide_structure: float
    avg_chars_per_slide: float
    visual_elements_per_slide: float
    text_heavy: bool
    lecture_style: bool
    generic_template: bool
    marketing_only: bool
    blurry: bool
    quote_collection: bool
    minimal_content: bool
    image_gallery: bool


def _pdf_page_text(page, *, allow_ocr: bool = True) -> str:
    text = (page.get_text() or "").strip()
    if len(text) >= 20 or not allow_ocr:
        return text
    try:
        tp = page.get_textpage_ocr()
        ocr_text = (page.get_text(textpage=tp) or "").strip()
        if len(ocr_text) > len(text):
            return ocr_text
    except Exception:
        pass
    return text


def _analyze_pptx(path: Path, thresholds: dict, *, fast: bool = False) -> QualityAnalysis:
    from pptx import Presentation

    prs = Presentation(str(path))
    slide_count = max(len(prs.slides), 1)
    total_chars = 0
    visual_elements = 0
    chart_table_count = 0
    title_only_slides = 0
    quote_slides = 0

    for slide in prs.slides:
        slide_chars = 0
        slide_visuals = 0
        slide_text_parts: list[str] = []
        for shape in slide.shapes:
            if shape.shape_type in VISUAL_SHAPE_TYPES:
                slide_visuals += 1
                visual_elements += 1
                if shape.shape_type in (MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.TABLE):
                    chart_table_count += 1
            if shape.has_text_frame:
                text = shape.text_frame.text or ""
                slide_text_parts.append(text.strip())
                slide_chars += len(text.strip())
        slide_text = " ".join(slide_text_parts)
        total_chars += slide_chars
        if slide_chars < 40 and slide_visuals == 0:
            title_only_slides += 1
        if '"' in slide_text and slide_chars < 220 and slide_visuals == 0:
            quote_slides += 1

    avg_chars = total_chars / slide_count
    visuals_per_slide = visual_elements / slide_count
    if fast:
        clarity_score, blurry = 70.0, False
    else:
        clarity_score, blurry = analyze_clarity(path)
    return _score_metrics(
        slide_count=slide_count,
        avg_chars=avg_chars,
        visuals_per_slide=visuals_per_slide,
        title_only_ratio=title_only_slides / slide_count,
        quote_ratio=quote_slides / slide_count,
        chart_table_count=chart_table_count,
        clarity_score=clarity_score,
        blurry=blurry,
        thresholds=thresholds,
    )


def _analyze_pdf(path: Path, thresholds: dict, *, fast: bool = False) -> QualityAnalysis:
    import fitz

    doc = fitz.open(str(path))
    slide_count = max(doc.page_count, 1)
    total_chars = 0
    visual_elements = 0
    title_only_pages = 0
    quote_pages = 0

    for page in doc:
        text = _pdf_page_text(page, allow_ocr=not fast)
        chars = len(text)
        images = len(page.get_images(full=True))
        total_chars += chars
        visual_elements += images
        if chars < 40 and images == 0:
            title_only_pages += 1
        if '"' in text and chars < 220 and images == 0:
            quote_pages += 1
    doc.close()

    avg_chars = total_chars / slide_count
    visuals_per_slide = visual_elements / slide_count
    if fast:
        clarity_score, blurry = 70.0, False
    else:
        clarity_score, blurry = analyze_clarity(path)
    return _score_metrics(
        slide_count=slide_count,
        avg_chars=avg_chars,
        visuals_per_slide=visuals_per_slide,
        title_only_ratio=title_only_pages / slide_count,
        quote_ratio=quote_pages / slide_count,
        chart_table_count=0,
        clarity_score=clarity_score,
        blurry=blurry,
        thresholds=thresholds,
    )


def _score_metrics(
    *,
    slide_count: int,
    avg_chars: float,
    visuals_per_slide: float,
    title_only_ratio: float,
    quote_ratio: float,
    chart_table_count: int,
    clarity_score: float,
    blurry: bool,
    thresholds: dict,
) -> QualityAnalysis:
    text_cfg = thresholds.get("text_density", {})
    gfx_cfg = thresholds.get("graphics_density", {})
    clarity_cfg = thresholds.get("visual_clarity", {})
    modern_cfg = thresholds.get("design_modernity", {})
    weights = thresholds.get("weights", {})

    max_chars = float(text_cfg.get("max_chars_per_slide", 800))
    lecture_threshold = float(text_cfg.get("lecture_mode_threshold", 0.80))
    text_ratio = min(1.0, avg_chars / max_chars) if max_chars else 0.0
    text_score = max(0.0, (1.0 - text_ratio) * 100.0)

    preferred_visuals = float(gfx_cfg.get("preferred_min", 1.5))
    min_visuals = float(gfx_cfg.get("min_visual_elements_per_slide", 0.5))
    visual_ratio = min(1.0, visuals_per_slide / preferred_visuals) if preferred_visuals else 0.0
    graphics_score = visual_ratio * 100.0

    clarity = clarity_score
    modernity = float(modern_cfg.get("min_score", 0.50)) * 100.0
    structure = min(100.0, slide_count * 8.0)

    if clarity_cfg.get("reject_if_pixelated", True) and blurry:
        clarity = min(clarity, 30.0)

    w_gfx = float(weights.get("graphics_density", 0.30))
    w_text = float(weights.get("text_density", 0.20))
    w_clarity = float(weights.get("visual_clarity", 0.20))
    w_modern = float(weights.get("design_modernity", 0.15))
    w_struct = float(weights.get("slide_structure", 0.15))

    quality = (
        graphics_score * w_gfx
        + text_score * w_text
        + clarity * w_clarity
        + modernity * w_modern
        + structure * w_struct
    )

    text_heavy = text_ratio > float(text_cfg.get("max_text_area_ratio", 0.65))
    lecture_style = text_ratio >= lecture_threshold
    generic_template = title_only_ratio >= 0.75 and visuals_per_slide < min_visuals
    marketing_only = avg_chars < 25 and visuals_per_slide < min_visuals and slide_count <= 12
    quote_collection = quote_ratio >= 0.5 and visuals_per_slide < 0.4
    minimal_content = avg_chars < 15 and slide_count >= 5
    image_gallery = (
        visuals_per_slide >= 2.0
        and avg_chars < 30
        and chart_table_count == 0
        and text_ratio < 0.2
    )

    return QualityAnalysis(
        quality=round(quality, 2),
        graphics_density=round(graphics_score, 2),
        text_density=round(text_ratio * 100.0, 2),
        clarity=round(clarity, 2),
        modernity=round(modernity, 2),
        slide_structure=round(structure, 2),
        avg_chars_per_slide=round(avg_chars, 1),
        visual_elements_per_slide=round(visuals_per_slide, 2),
        text_heavy=text_heavy,
        lecture_style=lecture_style,
        generic_template=generic_template,
        marketing_only=marketing_only,
        blurry=blurry,
        quote_collection=quote_collection,
        minimal_content=minimal_content,
        image_gallery=image_gallery,
    )


def analyze_quality(
    path: Path,
    thresholds: dict | None = None,
    *,
    fast: bool = False,
) -> QualityAnalysis:
    """Return quality metrics for a PPTX or PDF file.

    fast=True skips OCR and CV clarity (required for multi-million throughput).
    """
    thresholds = thresholds or {}
    ext = path.suffix.lower()
    if ext == ".pptx":
        return _analyze_pptx(path, thresholds, fast=fast)
    if ext == ".pdf":
        return _analyze_pdf(path, thresholds, fast=fast)
    raise ValueError(f"Unsupported format for quality analysis: {ext}")
