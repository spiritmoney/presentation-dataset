"""Computer-vision checks for blur and scan quality."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

import cv2
import numpy as np
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _laplacian_variance(gray: np.ndarray) -> float:
    return float(cv2.Laplacian(gray, cv2.CV_64F).var())


def _sharpness_score(variance: float) -> float:
    # Typical sharp slides: variance > 100; blurry scans: < 50
    return min(100.0, max(0.0, variance / 2.0))


def analyze_pdf_clarity(path: Path, min_variance: float = 50.0) -> tuple[float, bool]:
    import fitz

    doc = fitz.open(str(path))
    if doc.page_count == 0:
        doc.close()
        return 55.0, False
    page = doc[0]
    pix = page.get_pixmap(matrix=fitz.Matrix(150 / 72, 150 / 72))
    doc.close()
    arr = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
    if pix.n == 4:
        arr = cv2.cvtColor(arr, cv2.COLOR_RGBA2GRAY)
    elif pix.n == 3:
        arr = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
    else:
        arr = arr[:, :, 0]
    variance = _laplacian_variance(arr)
    score = _sharpness_score(variance)
    return score, variance < min_variance


def analyze_pptx_clarity(path: Path, min_variance: float = 50.0) -> tuple[float, bool]:
    from PIL import Image
    from pptx import Presentation

    prs = Presentation(str(path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = Image.open(BytesIO(shape.image.blob)).convert("L")
                    arr = np.array(img)
                    variance = _laplacian_variance(arr)
                    score = _sharpness_score(variance)
                    return score, variance < min_variance
                except Exception:
                    continue
    return 55.0, False


def analyze_clarity(path: Path, min_variance: float = 50.0) -> tuple[float, bool]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return analyze_pdf_clarity(path, min_variance)
    if ext == ".pptx":
        return analyze_pptx_clarity(path, min_variance)
    return 55.0, False
