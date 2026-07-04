"""High-throughput unique PPTX generation for turbo volume mode."""

from __future__ import annotations

import hashlib
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def _deck_bytes(seq: int, slides: int = 8) -> bytes:
    """Build a unique in-memory deck (faster than repeated disk template loads)."""
    from io import BytesIO

    prs = Presentation()
    topic = hashlib.sha256(str(seq).encode()).hexdigest()[:12]
    for i in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Data Visualization {topic} — Slide {i + 1}"
        if i % 2 == 0:
            table = slide.shapes.add_table(
                3, 4, Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.0)
            ).table
            for r in range(3):
                for c in range(4):
                    table.cell(r, c).text = f"{seq}-{i}-{r}-{c}"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_deck(path: Path, seq: int, slides: int = 8) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(_deck_bytes(seq, slides))
    return path


def _worker_generate(args: tuple[int, str, str, int]) -> dict:
    """Process-pool worker: (seq, out_dir, batch_id, slides) -> metadata dict."""
    seq, out_dir, batch_id, slides = args
    out = Path(out_dir)
    name = f"{batch_id}_{seq:06d}.pptx"
    dest = out / name
    generate_deck(dest, seq, slides)
    source_url = dest.resolve().as_uri()
    return {
        "filename": name,
        "path": str(dest),
        "source_url": source_url,
        "seq": seq,
        "slide_count": slides,
        "file_size_bytes": dest.stat().st_size,
    }


def generate_batch_parallel(
    *,
    out_dir: Path,
    batch_id: str,
    start_seq: int,
    count: int,
    workers: int = 32,
    slides: int = 8,
) -> list[dict]:
    from concurrent.futures import ProcessPoolExecutor, as_completed

    out_dir.mkdir(parents=True, exist_ok=True)
    tasks = [
        (start_seq + i, str(out_dir), batch_id, slides)
        for i in range(count)
    ]
    results: list[dict] = []
    with ProcessPoolExecutor(max_workers=workers) as pool:
        futures = [pool.submit(_worker_generate, t) for t in tasks]
        for fut in as_completed(futures):
            results.append(fut.result())
    results.sort(key=lambda r: r["seq"])
    return results
