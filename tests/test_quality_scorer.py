"""Tests for automated quality scoring."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from src.analysis.quality_scorer import analyze_quality


def _deck_with_tables(path: Path, slides: int = 6) -> None:
    prs = Presentation()
    for i in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Data {i}"
        slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(5), Inches(3))
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def test_quality_scorer_accepts_visual_deck(tmp_path: Path):
    path = tmp_path / "visual.pptx"
    _deck_with_tables(path, 8)
    analysis = analyze_quality(path)
    assert analysis.quality >= 60
    assert not analysis.text_heavy
    assert not analysis.lecture_style


def test_quality_scorer_flags_lecture_style(tmp_path: Path):
    prs = Presentation()
    for i in range(6):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        body = slide.placeholders[1]
        body.text = "Lorem ipsum " * 120
    path = tmp_path / "lecture.pptx"
    prs.save(str(path))
    analysis = analyze_quality(path)
    assert analysis.lecture_style or analysis.text_heavy
