"""Tests for automatic collection export."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from src.delivery.finalize import export_collection_artifacts


def _make_pptx(path: Path, slides: int = 5) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i in range(slides):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.text = f"Slide {i + 1}"
    prs.save(str(path))


def test_export_collection_artifacts_writes_under_data(tmp_path: Path):
    qualified = tmp_path / "qualified" / "BATCH-001"
    qualified.mkdir(parents=True)
    _make_pptx(qualified / "BATCH-001_000001.pptx", slides=5)

    manifests = tmp_path / "manifests"
    manifests.mkdir()
    (manifests / "BATCH-001.csv").write_text(
        "filename,file_type,source_url,download_domain,download_url,original_filename,"
        "collection_timestamp,download_timestamp,publication_date,author,organization,"
        "document_title,language,file_size_bytes,slide_count,quality_score,batch_id,"
        "source_status,public_access_status,tags,crawl_metadata,processing_metadata\n"
        "BATCH-001_000001.pptx,pptx,https://example.com/a.pptx,example.com,"
        "https://example.com/a.pptx,a.pptx,2026-01-01T00:00:00Z,2026-01-01T00:00:00Z,,"
        ",,Title,en,100,5,70,BATCH-001,reachable,PASS,[],{},{}",
        encoding="utf-8",
    )

    result = export_collection_artifacts(
        tmp_path,
        target_count=1,
        accepted_count=1,
        batch_id="BATCH-001",
    )

    assert result.report_path.parent == tmp_path / "reports"
    assert result.report_path.exists()
    assert (tmp_path / "reports" / "progress_latest.json").exists()
    assert result.manifest_csv == tmp_path / "manifests" / "MASTER_MANIFEST.csv"
    assert result.manifest_csv.exists()
    if result.manifest_xlsx is not None:
        assert result.manifest_xlsx.exists()
    assert result.delivery_zip.parent == tmp_path / "delivery"
    assert result.delivery_zip.exists()

    latest = (tmp_path / "reports" / "progress_latest.json").read_text(encoding="utf-8")
    assert "goal_reached" in latest
