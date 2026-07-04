"""Integration test for full pipeline stages (web sources only)."""

from __future__ import annotations

import json
import threading
from http.server import SimpleHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from src.pipeline.handlers import (
    StageContext,
    stage_dedupe,
    stage_download,
    stage_filter,
    stage_package,
    stage_report,
    stage_score,
    stage_validate,
)


def _make_pptx(path: Path, slides: int = 6) -> None:
    prs = Presentation()
    for i in range(slides):
        layout = prs.slide_layouts[1] if i > 0 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Analytics Slide {i + 1}"
        if i % 2 == 0:
            table = slide.shapes.add_table(3, 3, Inches(1), Inches(2), Inches(6), Inches(3)).table
            table.cell(0, 0).text = "Metric"
            table.cell(0, 1).text = "Q1"
            table.cell(0, 2).text = "Q2"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


@pytest.fixture
def pipeline_env(tmp_path: Path):
    seeds = tmp_path / "web_fixtures"
    seeds.mkdir(parents=True, exist_ok=True)
    _make_pptx(seeds / "deck1.pptx", 8)
    _make_pptx(seeds / "deck2.pptx", 7)

    class Handler(SimpleHTTPRequestHandler):
        def __init__(self, *args, **kwargs):
            super().__init__(*args, directory=str(seeds), **kwargs)

        def log_message(self, format, *args):  # noqa: A003
            return

    server = ThreadingHTTPServer(("127.0.0.1", 0), Handler)
    port = server.server_address[1]
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()

    batch_id = "BATCH-TEST-001"
    urls_dir = tmp_path / "staging" / "urls"
    urls_dir.mkdir(parents=True, exist_ok=True)
    queue = urls_dir / f"{batch_id}.jsonl"
    for fp in seeds.glob("*.pptx"):
        row = {
            "url": f"http://127.0.0.1:{port}/{fp.name}",
            "source_query": "test_http",
            "category": "web",
            "discovered_at": "2026-01-01T00:00:00+00:00",
        }
        with queue.open("a", encoding="utf-8") as f:
            f.write(json.dumps(row) + "\n")

    yield tmp_path, batch_id
    server.shutdown()


def test_full_pipeline_stages(pipeline_env, monkeypatch):
    data_dir, batch_id = pipeline_env
    ctx = StageContext(batch_id, data_dir)

    # Integration fixtures are minimal decks — focus this test on web-source wiring.
    class _Analysis:
        quality = 75.0
        graphics_density = 70.0
        text_density = 60.0
        clarity = 70.0
        modernity = 65.0
        slide_structure = 70.0
        avg_chars_per_slide = 80.0
        visual_elements_per_slide = 3.0
        blurry = False
        lecture_style = False
        text_heavy = False
        quote_collection = False
        minimal_content = False
        image_gallery = False
        generic_template = False
        marketing_only = False

    monkeypatch.setattr("src.pipeline.handlers.analyze_quality", lambda *a, **k: _Analysis())

    stage_download(ctx)
    stage_validate(ctx)
    stage_filter(ctx)
    stage_score(ctx)
    stage_dedupe(ctx)
    stage_package(ctx)
    stage_report(ctx)

    qualified = list((data_dir / "qualified" / batch_id).glob("*.pptx"))
    assert len(qualified) >= 2

    manifest = data_dir / "manifests" / f"{batch_id}.csv"
    assert manifest.exists()
    text = manifest.read_text(encoding="utf-8")
    assert "http://127.0.0.1" in text
    assert "file://" not in text

    audit = data_dir / "audit" / f"{batch_id}.jsonl"
    assert audit.exists()
