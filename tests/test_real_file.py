"""Tests for real-file and web-source qualification gates."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from src.validation.real_file import is_real_presentation, purge_non_real_qualified
from src.validation.source_url import is_local_source, is_web_source_url, queue_row_is_web


def _make_pptx(path: Path, slides: int = 5) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i in range(slides):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.text = f"Slide {i + 1}"
    prs.save(str(path))


def test_is_web_source_url():
    assert is_web_source_url("https://example.com/a.pptx") is True
    assert is_web_source_url("http://example.com/a.pdf") is True
    assert is_web_source_url("file:///C:/seeds/a.pptx") is False
    assert is_web_source_url("") is False
    assert is_web_source_url(None) is False


def test_rejects_local_sources():
    local, reason = is_local_source(source_url="file:///C:/seeds/a.pptx")
    assert local is True
    assert "non-web" in reason

    local, _ = is_local_source(
        source_url="https://example.com/a.pptx",
        local_path="C:/seeds/a.pptx",
    )
    assert local is True

    local, _ = is_local_source(
        source_url="https://example.com/a.pptx",
        category="seed",
    )
    assert local is True

    local, _ = is_local_source(
        source_url="https://example.com/a.pptx",
        content_type="application/local",
    )
    assert local is True

    local, reason = is_local_source(source_url="https://example.com/a.pptx")
    assert local is False
    assert reason == ""


def test_queue_row_rejects_local():
    assert queue_row_is_web({"url": "https://example.com/a.pptx"}) is True
    assert queue_row_is_web({"url": "file:///C:/a.pptx"}) is False
    assert (
        queue_row_is_web({"url": "https://example.com/a.pptx", "local_path": "C:/a.pptx"})
        is False
    )
    assert (
        queue_row_is_web({"url": "https://example.com/a.pptx", "category": "seed"}) is False
    )


def test_is_real_presentation_accepts_valid_pptx(tmp_path: Path):
    path = tmp_path / "deck.pptx"
    _make_pptx(path, slides=6)
    ok, reason = is_real_presentation(path)
    assert ok is True
    assert reason == ""


def test_is_real_presentation_rejects_stub(tmp_path: Path):
    path = tmp_path / "stub.pptx"
    path.write_bytes(b"not a real pptx")
    ok, reason = is_real_presentation(path)
    assert ok is False
    assert reason


def test_is_real_presentation_rejects_too_few_slides(tmp_path: Path):
    path = tmp_path / "short.pptx"
    _make_pptx(path, slides=2)
    ok, reason = is_real_presentation(path, min_slides=5)
    assert ok is False
    assert "below minimum" in reason


def test_purge_non_real_qualified(tmp_path: Path):
    qualified = tmp_path / "qualified" / "BATCH-1"
    qualified.mkdir(parents=True)
    good = qualified / "good.pptx"
    bad = qualified / "bad.pptx"
    _make_pptx(good, slides=5)
    bad.write_bytes(b"x")

    stats = purge_non_real_qualified(tmp_path, web_only=False)
    assert stats["kept"] == 1
    assert stats["removed"] == 1
    assert good.exists()
    assert not bad.exists()


def test_purge_removes_non_web_sources(tmp_path: Path):
    qualified = tmp_path / "qualified" / "BATCH-1"
    qualified.mkdir(parents=True)
    web = qualified / "web.pptx"
    seed = qualified / "seed.pptx"
    _make_pptx(web, slides=5)
    _make_pptx(seed, slides=5)

    manifests = tmp_path / "manifests"
    manifests.mkdir()
    (manifests / "BATCH-1.csv").write_text(
        "filename,source_url\n"
        "web.pptx,https://example.com/web.pptx\n"
        "seed.pptx,file:///C:/seeds/seed.pptx\n",
        encoding="utf-8",
    )

    stats = purge_non_real_qualified(tmp_path, web_only=True)
    assert stats["kept"] == 1
    assert stats["removed"] == 1
    assert web.exists()
    assert not seed.exists()
