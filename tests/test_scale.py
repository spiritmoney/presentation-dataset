"""Scale-path tests for 6M / 12h collection."""

from pathlib import Path

from src.crawlers.url_queue import append_urls, claim_batch, pending_count
from src.deduplication.index import DedupeIndex
from src.supervisor.counter import increment_qualified_count, read_qualified_count


def test_qualified_counter_o1(tmp_path: Path):
    assert read_qualified_count(tmp_path) == 0
    assert increment_qualified_count(tmp_path, 100) == 100
    assert increment_qualified_count(tmp_path, 50) == 150
    assert read_qualified_count(tmp_path) == 150


def test_url_queue_streaming_claim(tmp_path: Path):
    urls = tmp_path / "urls"
    rows = [{"url": f"https://example.com/{i}.pdf"} for i in range(25)]
    assert append_urls(urls, rows) == 25
    assert pending_count(urls) == 25

    claimed = claim_batch(urls, "BATCH-1", 10)
    assert len(claimed) == 10
    assert pending_count(urls) == 15
    assert (urls / "BATCH-1.jsonl").exists()


def test_dedupe_content_hash_o1(tmp_path: Path):
    from pptx import Presentation
    from pptx.util import Inches

    def make(path: Path):
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        for _ in range(5):
            prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(path))

    a = tmp_path / "a.pptx"
    b = tmp_path / "b.pptx"
    make(a)
    b.write_bytes(a.read_bytes())

    index = DedupeIndex(tmp_path / "dedupe.json")
    assert index.find_duplicate(a) == (None, None)
    index.register(a, "a.pptx")
    match, name = index.find_duplicate(b)
    assert match == "content"
    assert name == "a.pptx"
    index.save()
    reloaded = DedupeIndex(tmp_path / "dedupe.json")
    match, name = reloaded.find_duplicate(b)
    assert match == "content"
